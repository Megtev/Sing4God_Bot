from pptx import Presentation

def words_prs(prs):
	text_runs = []

	for slide in prs.slides:
		for shape in slide.shapes:
			str_t = ''
			if not shape.has_text_frame:
				continue
			for paragraph in shape.text_frame.paragraphs:
				for run in paragraph.runs:
					if len(run.text) in (0,1) and run.text == ' ':
						str_t += str_t.rstrip() + run.text
					elif len(run.text) in (0,1):
						str_t += run.text
					elif run.text.startswith(',') and run.text.endswith(' '):
						str_t = str_t.rstrip() + run.text
					elif run.text.startswith(',') and not run.text.endswith(' '):
						str_t = str_t.rstrip() + run.text + '\n'
					elif run.text.endswith(' '):
						str_t += run.text
					else:
						str_t += run.text + '\n'
			text_runs.append(str_t)
	return text_runs